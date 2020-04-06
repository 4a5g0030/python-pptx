import os
from pptx import Presentation
from  pptx.util import Inches


IM_DIR = os.path.join("hconcat")

if __name__ == '__main__':
    images = os.listdir(IM_DIR)
    prs = Presentation()

    for file in images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = file
        pic = slide.shapes.add_picture(os.path.join(IM_DIR, file), top=Inches(2.47), left=Inches(0.5), height=Inches(2.56))
    prs.save('test.pptx')